#!/usr/bin/env python3
"""
Spidey Context Budget Estimator.

Estimates extracted text tokens from prompts and uploaded files before
submitting a turn.

This is an approximation rather than an exact tokenizer. It is designed
to identify large token contributors and estimate overall context usage.
"""

import argparse
import csv
import json
import re
import xml.etree.ElementTree as ET
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable, List, Optional

from bs4 import BeautifulSoup
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".xml",
}

DEFAULT_CONTEXT_LIMIT = 100_000


@dataclass
class FileEstimate:
    file: str
    extension: str
    size_mb: float
    characters: int
    words: int
    estimated_tokens: int
    risk_band: str
    extraction_notes: str


def clean_text(text: str) -> str:
    """Normalize extracted text before estimating tokens."""
    if not text:
        return ""

    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def estimate_tokens(text: str) -> int:
    """
    Estimate the number of tokens in extracted text.

    Uses tiktoken when available. Otherwise, it falls back to a
    conservative character- and word-based heuristic.
    """
    text = text or ""

    try:
        import tiktoken  # type: ignore

        encoder = tiktoken.get_encoding("cl100k_base")
        return len(encoder.encode(text))

    except Exception:
        characters = len(text)
        words = len(re.findall(r"\S+", text))

        return int(
            max(
                characters / 4.0,
                words * 1.35,
            )
        )


def risk_band(
    tokens: int,
    total_limit: int = DEFAULT_CONTEXT_LIMIT,
) -> str:
    """Assign a risk label based on the selected context limit."""
    if total_limit <= 0:
        raise ValueError(
            "The context limit must be greater than zero."
        )

    percentage = tokens / total_limit

    if percentage >= 0.95:
        return "CRITICAL"

    if percentage >= 0.85:
        return "LIKELY_OVER_LIMIT_SOON"

    if percentage >= 0.70:
        return "CLOSE_TO_LIMIT"

    if percentage >= 0.40:
        return "WATCH"

    return "SAFE"


def extract_pdf(
    path: Path,
    max_pages: Optional[int] = None,
) -> tuple[str, str]:
    """Extract text from a PDF."""
    try:
        import fitz
    except Exception as exc:
        return (
            "",
            f"PDF skipped: PyMuPDF not installed ({exc})",
        )

    notes = []
    chunks = []

    try:
        with fitz.open(path) as document:
            page_count = len(document)

            if max_pages is None:
                pages_to_read = page_count
            else:
                pages_to_read = min(
                    page_count,
                    max_pages,
                )

            for page_index in range(pages_to_read):
                page = document.load_page(page_index)
                chunks.append(
                    page.get_text("text")
                )

            notes.append(
                f"PDF pages read: "
                f"{pages_to_read}/{page_count}"
            )

            if (
                max_pages is not None
                and page_count > max_pages
            ):
                notes.append(
                    "PDF truncated by max-pages option"
                )

    except Exception as exc:
        return (
            "",
            f"PDF extraction failed: {exc}",
        )

    return (
        clean_text("\n".join(chunks)),
        "; ".join(notes),
    )


def extract_docx(path: Path) -> tuple[str, str]:
    """Extract paragraphs and tables from a DOCX file."""
    try:
        import docx
    except Exception as exc:
        return (
            "",
            f"DOCX skipped: python-docx not installed ({exc})",
        )

    chunks = []

    try:
        document = docx.Document(str(path))

        for paragraph in document.paragraphs:
            if paragraph.text:
                chunks.append(paragraph.text)

        for table in document.tables:
            for row in table.rows:
                chunks.append(
                    "\t".join(
                        cell.text
                        for cell in row.cells
                    )
                )

        notes = (
            f"DOCX paragraphs: "
            f"{len(document.paragraphs)}, "
            f"tables: {len(document.tables)}"
        )

        return (
            clean_text("\n".join(chunks)),
            notes,
        )

    except Exception as exc:
        return (
            "",
            f"DOCX extraction failed: {exc}",
        )


def extract_pptx_text(path: Path) -> str:
    """Extract slide text and speaker notes from a PPTX file."""
    presentation = Presentation(str(path))
    parts = []

    for slide_index, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        parts.append(
            f"\n--- Slide {slide_index} ---\n"
        )

        for shape in slide.shapes:
            if (
                hasattr(shape, "text")
                and shape.text
            ):
                parts.append(shape.text)

        if slide.has_notes_slide:
            notes_text_frame = (
                slide.notes_slide.notes_text_frame
            )

            if (
                notes_text_frame
                and notes_text_frame.text
            ):
                parts.append(
                    "\n[Speaker notes]\n"
                    f"{notes_text_frame.text}"
                )

    return "\n".join(parts)


def extract_xlsx(
    path: Path,
    max_cells: Optional[int] = None,
) -> tuple[str, str]:
    """Extract displayed cell values and formulas from an XLSX file."""
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        return (
            "",
            f"XLSX skipped: openpyxl not installed ({exc})",
        )

    chunks = []
    cell_count = 0

    try:
        workbook = load_workbook(
            filename=path,
            read_only=True,
            data_only=False,
        )

        sheet_count = len(
            workbook.sheetnames
        )

        for worksheet in workbook.worksheets:
            chunks.append(
                f"\n[SHEET: {worksheet.title}]\n"
            )

            for row in worksheet.iter_rows():
                values = []

                for cell in row:
                    if cell.value is not None:
                        values.append(
                            str(cell.value)
                        )

                        cell_count += 1

                        if (
                            max_cells is not None
                            and cell_count >= max_cells
                        ):
                            if values:
                                chunks.append(
                                    "\t".join(values)
                                )

                            notes = (
                                f"XLSX sheets: {sheet_count}; "
                                f"cells read: {cell_count}; "
                                "truncated by max-cells option"
                            )

                            workbook.close()

                            return (
                                clean_text(
                                    "\n".join(chunks)
                                ),
                                notes,
                            )

                if values:
                    chunks.append(
                        "\t".join(values)
                    )

        workbook.close()

        notes = (
            f"XLSX sheets: {sheet_count}; "
            f"non-empty cells read: {cell_count}"
        )

        return (
            clean_text("\n".join(chunks)),
            notes,
        )

    except Exception as exc:
        return (
            "",
            f"XLSX extraction failed: {exc}",
        )


def extract_csv(
    path: Path,
    max_rows: Optional[int] = None,
) -> tuple[str, str]:
    """Extract rows from a CSV or delimited-text file."""
    chunks = []
    row_count = 0

    try:
        with path.open(
            "r",
            encoding="utf-8-sig",
            errors="replace",
            newline="",
        ) as file:
            sample = file.read(4096)
            file.seek(0)

            try:
                dialect = csv.Sniffer().sniff(
                    sample
                )
            except Exception:
                dialect = csv.excel

            reader = csv.reader(
                file,
                dialect,
            )

            for row in reader:
                row_count += 1
                chunks.append(
                    "\t".join(row)
                )

                if (
                    max_rows is not None
                    and row_count >= max_rows
                ):
                    notes = (
                        f"CSV rows read: {row_count}; "
                        "truncated by max-rows option"
                    )

                    return (
                        clean_text(
                            "\n".join(chunks)
                        ),
                        notes,
                    )

        return (
            clean_text("\n".join(chunks)),
            f"CSV rows read: {row_count}",
        )

    except Exception as exc:
        return (
            "",
            f"CSV extraction failed: {exc}",
        )


def extract_html_text(path: Path) -> str:
    """Extract visible text from an HTML file."""
    raw = path.read_text(
        encoding="utf-8",
        errors="ignore",
    )

    soup = BeautifulSoup(
        raw,
        "html.parser",
    )

    for tag in soup(
        ["script", "style", "noscript"]
    ):
        tag.decompose()

    return soup.get_text(
        separator="\n",
        strip=True,
    )


def extract_xml_text(path: Path) -> str:
    """Extract tags and text values from an XML file."""
    raw = path.read_text(
        encoding="utf-8",
        errors="ignore",
    )

    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        return raw

    parts = []

    def walk(node, depth: int = 0) -> None:
        tag = node.tag

        if tag:
            parts.append(
                f"{'  ' * depth}<{tag}>"
            )

        if (
            node.text
            and node.text.strip()
        ):
            parts.append(
                node.text.strip()
            )

        for child in node:
            walk(
                child,
                depth + 1,
            )

        if (
            node.tail
            and node.tail.strip()
        ):
            parts.append(
                node.tail.strip()
            )

    walk(root)

    return "\n".join(parts)


def extract_text(
    path: Path,
) -> tuple[str, str]:
    """Route a file to the appropriate extraction function."""
    extension = path.suffix.lower()

    if extension == ".pdf":
        return extract_pdf(path)

    if extension == ".docx":
        return extract_docx(path)

    if extension == ".pptx":
        try:
            text = extract_pptx_text(path)

            return (
                clean_text(text),
                "PPTX slides and speaker notes read",
            )

        except Exception as exc:
            return (
                "",
                f"PPTX extraction failed: {exc}",
            )

    if extension == ".xlsx":
        return extract_xlsx(path)

    if extension == ".csv":
        return extract_csv(path)

    if extension in {".txt", ".md"}:
        try:
            text = path.read_text(
                encoding="utf-8",
                errors="replace",
            )

            return (
                clean_text(text),
                "Plain text/Markdown read",
            )

        except Exception as exc:
            return (
                "",
                f"Text extraction failed: {exc}",
            )

    if extension in {".html", ".htm"}:
        try:
            text = extract_html_text(path)

            return (
                clean_text(text),
                "HTML text extracted with scripts/styles removed",
            )

        except Exception as exc:
            return (
                "",
                f"HTML extraction failed: {exc}",
            )

    if extension == ".xml":
        try:
            text = extract_xml_text(path)

            return (
                clean_text(text),
                "XML text extracted from element tree",
            )

        except Exception as exc:
            return (
                "",
                f"XML extraction failed: {exc}",
            )

    return (
        "",
        f"Unsupported extension: {extension}",
    )


def iter_files(
    paths: Iterable[Path],
) -> Iterable[Path]:
    """Yield supported files from individual paths or folders."""
    for path in paths:
        if path.is_dir():
            for candidate in sorted(
                path.rglob("*")
            ):
                if (
                    candidate.is_file()
                    and candidate.suffix.lower()
                    in SUPPORTED_EXTENSIONS
                ):
                    yield candidate

        elif path.is_file():
            if path.suffix.lower() == ".zip":
                continue

            if (
                path.suffix.lower()
                in SUPPORTED_EXTENSIONS
            ):
                yield path


def estimate_file(
    path: Path,
    context_limit: int = DEFAULT_CONTEXT_LIMIT,
) -> FileEstimate:
    """Extract and estimate the token usage of one file."""
    text, notes = extract_text(path)

    characters = len(text)
    words = len(
        re.findall(r"\S+", text)
    )
    tokens = estimate_tokens(text)

    try:
        size_mb = (
            path.stat().st_size
            / (1024 * 1024)
        )
    except Exception:
        size_mb = 0.0

    return FileEstimate(
        file=str(path),
        extension=path.suffix.lower(),
        size_mb=round(size_mb, 2),
        characters=characters,
        words=words,
        estimated_tokens=tokens,
        risk_band=risk_band(
            tokens,
            context_limit,
        ),
        extraction_notes=notes,
    )


def estimate_prompt(
    prompt: str,
    context_limit: int = DEFAULT_CONTEXT_LIMIT,
) -> FileEstimate:
    """Estimate the token usage of manually entered prompt text."""
    text = prompt or ""

    characters = len(text)
    words = len(
        re.findall(r"\S+", text)
    )
    tokens = estimate_tokens(text)

    return FileEstimate(
        file="[PROMPT_TEXT]",
        extension="prompt",
        size_mb=0.0,
        characters=characters,
        words=words,
        estimated_tokens=tokens,
        risk_band=risk_band(
            tokens,
            context_limit,
        ),
        extraction_notes="Prompt text supplied manually",
    )


def summarize(
    estimates: List[FileEstimate],
    context_limit: int = DEFAULT_CONTEXT_LIMIT,
) -> dict:
    """Summarize total estimated context usage."""
    if context_limit <= 0:
        raise ValueError(
            "The context limit must be greater than zero."
        )

    total = sum(
        estimate.estimated_tokens
        for estimate in estimates
    )

    percentage = (
        total
        / context_limit
        * 100
    )

    if total >= context_limit:
        status = "OVER_LIMIT"

    elif percentage >= 95:
        status = "CRITICAL"

    elif percentage >= 85:
        status = "LIKELY_OVER_LIMIT_SOON"

    elif percentage >= 70:
        status = "CLOSE_TO_LIMIT"

    else:
        status = "SAFE"

    largest = sorted(
        estimates,
        key=lambda estimate: (
            estimate.estimated_tokens
        ),
        reverse=True,
    )[:5]

    return {
        "context_limit": context_limit,
        "total_estimated_tokens": total,
        "percent_of_limit": round(
            percentage,
            1,
        ),
        "status": status,
        "largest_contributors": [
            asdict(estimate)
            for estimate in largest
        ],
    }


def print_table(
    estimates: List[FileEstimate],
    context_limit: int,
) -> None:
    """Print estimates as a terminal table."""
    estimates_sorted = sorted(
        estimates,
        key=lambda estimate: (
            estimate.estimated_tokens
        ),
        reverse=True,
    )

    print(
        "\nEstimated extracted-token budget"
    )
    print("=" * 95)

    print(
        f"{'Tokens':>10}  "
        f"{'%Limit':>7}  "
        f"{'Words':>9}  "
        f"{'MB':>7}  "
        f"{'Risk':<22}  "
        "File"
    )

    print("-" * 95)

    for estimate in estimates_sorted:
        percentage = (
            estimate.estimated_tokens
            / context_limit
            * 100
        )

        print(
            f"{estimate.estimated_tokens:>10,}  "
            f"{percentage:>6.1f}%  "
            f"{estimate.words:>9,}  "
            f"{estimate.size_mb:>7.2f}  "
            f"{estimate.risk_band:<22}  "
            f"{Path(estimate.file).name}"
        )

    print("-" * 95)

    summary = summarize(
        estimates,
        context_limit,
    )

    print(
        f"TOTAL: "
        f"{summary['total_estimated_tokens']:,} / "
        f"{context_limit:,} tokens "
        f"({summary['percent_of_limit']}%)"
    )

    print(
        f"STATUS: {summary['status']}"
    )

    print(
        "\nNote: This estimates parsed/extracted text tokens "
        "rather than file size. The exact tokenizer used by "
        "the target model may differ."
    )


def main() -> None:
    """Run the estimator from the command line."""
    parser = argparse.ArgumentParser(
        description=(
            "Estimate extracted token counts for prompts "
            "and supported files."
        )
    )

    parser.add_argument(
        "paths",
        nargs="+",
        help=(
            "Files or folders to scan. Supported formats: "
            "PDF, DOCX, PPTX, XLSX, CSV, TXT, MD, HTML, "
            "and XML."
        ),
    )

    parser.add_argument(
        "--prompt-file",
        help=(
            "Optional prompt text file to include "
            "in the estimate."
        ),
    )

    parser.add_argument(
        "--limit",
        type=int,
        default=DEFAULT_CONTEXT_LIMIT,
        help=(
            "Context limit. "
            "Default: 100,000 tokens."
        ),
    )

    parser.add_argument(
        "--json-out",
        help=(
            "Optional output path for the generated "
            "JSON report."
        ),
    )

    args = parser.parse_args()

    if args.limit <= 0:
        parser.error(
            "--limit must be greater than zero."
        )

    estimates: List[FileEstimate] = []

    if args.prompt_file:
        prompt_path = Path(
            args.prompt_file
        )

        if not prompt_path.is_file():
            parser.error(
                f"Prompt file not found: "
                f"{prompt_path}"
            )

        prompt = prompt_path.read_text(
            encoding="utf-8",
            errors="replace",
        )

        estimates.append(
            estimate_prompt(
                prompt,
                args.limit,
            )
        )

    files = list(
        iter_files(
            Path(path)
            for path in args.paths
        )
    )

    for file_path in files:
        estimates.append(
            estimate_file(
                file_path,
                args.limit,
            )
        )

    if not estimates:
        print(
            "No supported files or prompt text were found."
        )
        return

    print_table(
        estimates,
        args.limit,
    )

    if args.json_out:
        report = {
            "summary": summarize(
                estimates,
                args.limit,
            ),
            "files": [
                asdict(estimate)
                for estimate in estimates
            ],
        }

        output_path = Path(
            args.json_out
        )

        output_path.write_text(
            json.dumps(
                report,
                indent=2,
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )

        print(
            f"\nJSON report written to: "
            f"{output_path}"
        )


if __name__ == "__main__":
    main()
